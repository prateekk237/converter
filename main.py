import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import pytesseract
import io
import os
import cv2
import numpy as np


# Helper Function: Extract Font Styles
def extract_font_styles(text_info):
    font_name = text_info.get("font", "Arial")
    font_size = text_info.get("size", 12)
    raw_color = text_info.get("color", 0)

    # Ensure color is a valid RGB tuple
    if isinstance(raw_color, tuple) and len(raw_color) == 3:
        color = raw_color
    else:
        color = (0, 0, 0)  # Default to black if invalid

    return font_name, font_size, color


# Helper Function: Detect Shapes in Images
def detect_shapes(image_path):
    image = cv2.imread(image_path)
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    _, thresh = cv2.threshold(gray, 127, 255, cv2.THRESH_BINARY)
    contours, _ = cv2.findContours(thresh, cv2.RETR_TREE, cv2.CHAIN_APPROX_SIMPLE)

    shapes = []
    for contour in contours:
        approx = cv2.approxPolyDP(contour, 0.02 * cv2.arcLength(contour, True), True)
        if len(approx) == 3:
            shapes.append(("Triangle", contour))
        elif len(approx) == 4:
            shapes.append(("Rectangle", contour))
        else:
            shapes.append(("Circle", contour))

    return shapes


# Main Function: PDF to PPTX Conversion
def pdf_to_pptx_advanced(pdf_file, pptx_file):
    presentation = Presentation()
    pdf_document = fitz.open(pdf_file)

    for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]
        text_instances = page.get_text("dict")["blocks"]
        images = page.get_images(full=True)

        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank layout

        # Process Text Blocks
        for block in text_instances:
            if block["type"] == 0:  # Text block
                for line in block["lines"]:
                    for span in line["spans"]:
                        font_name, font_size, color = extract_font_styles(span)
                        textbox = slide.shapes.add_textbox(
                            Inches(span["bbox"][0] / 72),
                            Inches(span["bbox"][1] / 72),
                            Inches((span["bbox"][2] - span["bbox"][0]) / 72),
                            Inches((span["bbox"][3] - span["bbox"][1]) / 72)
                        )
                        text_frame = textbox.text_frame
                        text_frame.text = span["text"]
                        for paragraph in text_frame.paragraphs:
                            paragraph.font.name = font_name
                            paragraph.font.size = Pt(font_size)
                            paragraph.font.color.rgb = RGBColor(*color)
                            paragraph.alignment = PP_ALIGN.LEFT

        # Process Images
        for img_index, img in enumerate(images):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"].upper()

            # Save Image to Temporary Path
            image_path = f"temp_image_{page_num}_{img_index}.{image_ext.lower()}"
            with open(image_path, "wb") as img_file:
                img_file.write(image_bytes)

            # Detect and Recreate Shapes in PowerPoint
            shapes = detect_shapes(image_path)
            for shape_type, contour in shapes:
                if shape_type == "Rectangle":
                    slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(1),
                        Inches(1),
                        Inches(2),
                        Inches(1)
                    )
                elif shape_type == "Triangle":
                    slide.shapes.add_shape(
                        MSO_SHAPE.ISOSCELES_TRIANGLE,
                        Inches(1),
                        Inches(1),
                        Inches(2),
                        Inches(2)
                    )
                elif shape_type == "Circle":
                    slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        Inches(1),
                        Inches(1),
                        Inches(2),
                        Inches(2)
                    )

            # Remove Temporary Image File
            os.remove(image_path)

    presentation.save(pptx_file)


# Streamlit App
import streamlit as st

st.title("Advanced PDF to PPTX Converter")
st.write("Upload a PDF to convert it into a PowerPoint presentation with editable text and shapes.")

uploaded_file = st.file_uploader("Choose a PDF file", type="pdf")

if uploaded_file is not None:
    # Save uploaded file temporarily
    pdf_path = f"uploaded_{uploaded_file.name}"
    with open(pdf_path, "wb") as f:
        f.write(uploaded_file.read())

    output_pptx_path = "converted_presentation.pptx"

    # Convert PDF to PPTX
    with st.spinner("Converting PDF to PPTX..."):
        pdf_to_pptx_advanced(pdf_path, output_pptx_path)

    # Provide download link for PPTX
    with open(output_pptx_path, "rb") as pptx_file:
        st.success("Conversion completed!")
        st.download_button(
            label="Download PPTX",
            data=pptx_file,
            file_name="converted_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    # Cleanup temporary files
    os.remove(pdf_path)
    os.remove(output_pptx_path)
